import logging
from pathlib import Path
from typing import List, Union

from langchain_core.documents import Document

logger = logging.getLogger(__name__)

SUPPORTED_EXTENSIONS = {
    ".pdf", ".txt", ".pptx", ".docx", ".jpeg", ".jpg", ".png",
    ".csv", ".xlsx", ".xls", ".db",
}


def _is_meaningful(text: str) -> bool:
    return bool(text and text.strip())


def _load_pdf(file_path: Path) -> List[Document]:
    from langchain_community.document_loaders import PyPDFLoader
    try:
        loader = PyPDFLoader(str(file_path))
        pages = loader.load()
    except Exception as exc:
        logger.error("Failed to load PDF '%s': %s", file_path.name, exc)
        return []
    docs = []
    for page in pages:
        if not _is_meaningful(page.page_content):
            continue
        page.metadata["source"] = file_path.name
        page.metadata["page"] = page.metadata.get("page", 0)
        page.metadata["page_display"] = page.metadata["page"] + 1
        docs.append(page)
    return docs


def _load_txt(file_path: Path) -> List[Document]:
    try:
        text = file_path.read_text(encoding="utf-8", errors="replace")
    except Exception as exc:
        logger.error("Failed to read TXT '%s': %s", file_path.name, exc)
        return []
    if not _is_meaningful(text):
        return []
    # Split on double newlines as rough pages
    sections = [s.strip() for s in text.split("\n\n") if s.strip()]
    docs = []
    for i, section in enumerate(sections):
        docs.append(Document(
            page_content=section,
            metadata={"source": file_path.name, "page": i, "page_display": i + 1}
        ))
    return docs


def _load_pptx(file_path: Path) -> List[Document]:
    try:
        from pptx import Presentation
    except ImportError:
        logger.error("python-pptx not installed. Cannot load .pptx files.")
        return []
    try:
        prs = Presentation(str(file_path))
    except Exception as exc:
        logger.error("Failed to open PPTX '%s': %s", file_path.name, exc)
        return []
    docs = []
    for i, slide in enumerate(prs.slides):
        texts = []
        for shape in slide.shapes:
            if hasattr(shape, "text") and shape.text.strip():
                texts.append(shape.text.strip())
        content = "\n".join(texts)
        if not _is_meaningful(content):
            continue
        docs.append(Document(
            page_content=content,
            metadata={"source": file_path.name, "page": i, "page_display": i + 1}
        ))
    return docs


def _load_docx(file_path: Path) -> List[Document]:
    try:
        import docx
    except ImportError:
        logger.error("python-docx not installed. Cannot load .docx files.")
        return []
    try:
        doc = docx.Document(str(file_path))
    except Exception as exc:
        logger.error("Failed to open DOCX '%s': %s", file_path.name, exc)
        return []

    # Collect all non-empty paragraph text
    paragraphs = [p.text.strip() for p in doc.paragraphs if p.text.strip()]

    # Also read tables
    for table in doc.tables:
        for row in table.rows:
            row_text = " | ".join(cell.text.strip() for cell in row.cells if cell.text.strip())
            if row_text:
                paragraphs.append(row_text)

    if not paragraphs:
        return []

    # Group paragraphs into ~500-character blocks
    BLOCK_TARGET = 500
    blocks: List[str] = []
    current_block: List[str] = []
    current_len = 0

    for para in paragraphs:
        current_block.append(para)
        current_len += len(para) + 1  # +1 for newline separator
        if current_len >= BLOCK_TARGET:
            blocks.append("\n".join(current_block))
            current_block = []
            current_len = 0

    if current_block:
        blocks.append("\n".join(current_block))

    docs = []
    for i, block in enumerate(blocks):
        if _is_meaningful(block):
            docs.append(Document(
                page_content=block,
                metadata={
                    "source": file_path.name,
                    "page": i,
                    "page_display": i + 1,
                    "file_type": "docx",
                },
            ))
    return docs


def _load_csv(file_path: Path) -> List[Document]:
    try:
        import pandas as pd
    except ImportError:
        logger.error("pandas not installed. Cannot load .csv files.")
        return []
    try:
        df = pd.read_csv(str(file_path))
    except Exception as exc:
        logger.error("Failed to read CSV '%s': %s", file_path.name, exc)
        return []

    MAX_ROWS = 5000
    if len(df) > MAX_ROWS:
        logger.warning(
            "CSV '%s' has %d rows; loading only the first %d.",
            file_path.name, len(df), MAX_ROWS,
        )
        df = df.head(MAX_ROWS)

    docs = []
    columns = list(df.columns)
    for row_idx, row in df.iterrows():
        pairs = [f"{col}: {row[col]}" for col in columns if pd.notna(row[col])]
        content = ", ".join(pairs)
        if _is_meaningful(content):
            docs.append(Document(
                page_content=content,
                metadata={
                    "source": file_path.name,
                    "page": int(row_idx),
                    "page_display": int(row_idx) + 1,
                    "file_type": "csv",
                },
            ))
    return docs


def _load_excel(file_path: Path) -> List[Document]:
    try:
        import pandas as pd
    except ImportError:
        logger.error("pandas not installed. Cannot load Excel files.")
        return []
    try:
        sheets = pd.read_excel(str(file_path), sheet_name=None)
    except Exception as exc:
        logger.error("Failed to read Excel '%s': %s", file_path.name, exc)
        return []

    MAX_ROWS = 5000
    docs = []
    for sheet_name, df in sheets.items():
        if len(df) > MAX_ROWS:
            logger.warning(
                "Excel '%s' sheet '%s' has %d rows; loading only the first %d.",
                file_path.name, sheet_name, len(df), MAX_ROWS,
            )
            df = df.head(MAX_ROWS)

        columns = list(df.columns)
        for row_idx, row in df.iterrows():
            pairs = [f"{col}: {row[col]}" for col in columns if pd.notna(row[col])]
            content = ", ".join(pairs)
            if _is_meaningful(content):
                docs.append(Document(
                    page_content=content,
                    metadata={
                        "source": file_path.name,
                        "sheet": str(sheet_name),
                        "page": int(row_idx),
                        "page_display": int(row_idx) + 1,
                        "file_type": "excel",
                    },
                ))
    return docs


def _load_image(file_path: Path) -> List[Document]:
    import base64
    from config import Config

    if not Config.OPENAI_API_KEY:
        logger.error("OPENAI_API_KEY is not set. Cannot transcribe image '%s'.", file_path.name)
        return []
    try:
        with open(file_path, "rb") as image_file:
            base64_image = base64.b64encode(image_file.read()).decode("utf-8")
    except Exception as exc:
        logger.error("Failed to read image file '%s': %s", file_path.name, exc)
        return []

    ext = file_path.suffix.lower()
    mime_type = "image/png" if ext == ".png" else "image/jpeg"

    try:
        from openai import OpenAI
        client = OpenAI(api_key=Config.OPENAI_API_KEY)
        response = client.chat.completions.create(
            model="gpt-4o",
            messages=[
                {
                    "role": "user",
                    "content": [
                        {
                            "type": "text",
                            "text": (
                                "Please transcribe all readable text from this image exactly. "
                                "If the image contains charts, tables, or diagrams, describe them in detail. "
                                "If it is a photograph or illustration, describe the scene in detail. "
                                "Do not include any conversational remarks or introductory text, just output the extracted content."
                            ),
                        },
                        {
                            "type": "image_url",
                            "image_url": {
                                "url": f"data:{mime_type};base64,{base64_image}"
                            },
                        },
                    ],
                }
            ],
            max_tokens=1500,
        )
        content = response.choices[0].message.content
        if not content or not content.strip():
            return []

        return [Document(
            page_content=content,
            metadata={"source": file_path.name, "page": 0, "page_display": 1}
        )]
    except Exception as exc:
        logger.error("Failed to transcribe image '%s' using OpenAI: %s", file_path.name, exc)
        return []


def load_document(file_path: Union[str, Path]) -> List[Document]:
    file_path = Path(file_path)
    if not file_path.exists():
        logger.warning("File not found: %s", file_path)
        return []
    if file_path.stat().st_size == 0:
        logger.warning("Empty file: %s", file_path)
        return []

    ext = file_path.suffix.lower()
    if ext == ".pdf":
        docs = _load_pdf(file_path)
    elif ext == ".txt":
        docs = _load_txt(file_path)
    elif ext == ".pptx":
        docs = _load_pptx(file_path)
    elif ext == ".docx":
        docs = _load_docx(file_path)
    elif ext == ".csv":
        docs = _load_csv(file_path)
    elif ext in {".xlsx", ".xls"}:
        docs = _load_excel(file_path)
    elif ext in {".jpeg", ".jpg", ".png"}:
        docs = _load_image(file_path)
    elif ext == ".db":
        # .db files are handled separately in app.py via SQLAgent
        logger.info("SQLite database '%s' detected — skipping document loader.", file_path.name)
        return []
    else:
        logger.warning("Unsupported file type '%s': %s", ext, file_path.name)
        return []

    if not docs:
        logger.warning("No extractable text in '%s'.", file_path.name)
    return docs


def load_documents_from_paths(file_paths: List[Union[str, Path]]) -> List[Document]:
    all_docs = []
    for path in file_paths:
        docs = load_document(path)
        all_docs.extend(docs)
    return all_docs


def load_documents_from_directory(directory: Union[str, Path]) -> List[Document]:
    directory = Path(directory)
    all_docs = []
    if not directory.exists():
        logger.warning("Directory does not exist: %s", directory)
        return all_docs
    files = sorted(
        f for f in directory.iterdir()
        if f.suffix.lower() in SUPPORTED_EXTENSIONS
    )
    for f in files:
        logger.info("Loading: %s", f.name)
        docs = load_document(f)
        logger.info("  -> %d section(s) loaded", len(docs))
        all_docs.extend(docs)
    return all_docs


# Backward compat aliases
load_pdf = load_document
load_pdfs_from_paths = load_documents_from_paths
load_pdfs_from_directory = load_documents_from_directory

